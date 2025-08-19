"""
All document chunker classes for text, code, tables, multimodal, and hybrid chunking.
"""
from typing import List, Callable, Optional, Any, Union, Dict
from dataclasses import dataclass
from enum import Enum
import re
import numpy as np
# Optional spacy import for NLP features
try:
    import spacy
    SPACY_AVAILABLE = True
except ImportError:
    SPACY_AVAILABLE = False
    print("Warning: spacy not available. NLP features will be disabled.")

# Optional transformers import for advanced document processing
try:
    from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
    _AUTO_MODEL_CLASS = AutoModelForSeq2SeqLM
    TRANSFORMERS_AVAILABLE = True
except ImportError:
    try:
        from transformers import AutoTokenizer, AutoModelForSeq2SeqGeneration
        _AUTO_MODEL_CLASS = AutoModelForSeq2SeqGeneration
        TRANSFORMERS_AVAILABLE = True
    except ImportError:
        try:
            from transformers import AutoTokenizer
            _AUTO_MODEL_CLASS = None
            TRANSFORMERS_AVAILABLE = True
        except ImportError:
            TRANSFORMERS_AVAILABLE = False
            _AUTO_MODEL_CLASS = None
            print("Warning: transformers not available. Advanced document processing features will be disabled.")

try:
    import nltk
    nltk.download('punkt', quiet=True)
    from nltk.tokenize import sent_tokenize
    _HAS_NLTK = True
except ImportError:
    _HAS_NLTK = False

class SemanticChunker:
    """Implements semantic document chunking."""
    def __init__(self, model, min_chunk_size: int = 100, max_chunk_size: int = 1000, similarity_threshold: float = 0.7, **kwargs):
        self.model = model
        self.min_chunk_size = min_chunk_size
        self.max_chunk_size = max_chunk_size
        self.similarity_threshold = similarity_threshold
        
        if TRANSFORMERS_AVAILABLE:
            self.tokenizer = AutoTokenizer.from_pretrained("facebook/bart-large-cnn")
            
            # Backward compatible model loading
            if _AUTO_MODEL_CLASS is not None:
                self.summarizer = _AUTO_MODEL_CLASS.from_pretrained("facebook/bart-large-cnn")
            else:
                # Fallback for very old versions - try to import the model directly
                try:
                    from transformers import BartForConditionalGeneration
                    self.summarizer = BartForConditionalGeneration.from_pretrained("facebook/bart-large-cnn")
                except ImportError:
                    raise ImportError("Unable to load BART model. Please ensure transformers is properly installed.")
        else:
            self.tokenizer = None
            self.summarizer = None
    async def chunk_document(self, text: str, metadata: Optional[Dict[str, Any]] = None, **kwargs) -> List[Any]:
        sentences = self._split_into_sentences(text)
        sentence_embeddings = await self.model.embeddings(sentences)
        chunks = self._group_similar_sentences(sentences, sentence_embeddings)
        return [
            {
                'text': chunk_text,
                'metadata': metadata or {},
                'chunk_id': f"chunk_{i}",
                'parent_id': None,
                'semantic_score': self._calculate_semantic_score(chunk_text)
            }
            for i, chunk_text in enumerate(chunks)
        ]
    def _split_into_sentences(self, text: str) -> List[str]:
        if SPACY_AVAILABLE:
            doc = spacy.load("en_core_web_sm")(text)
            return [sent.text.strip() for sent in doc.sents]
        else:
            # Fallback to simple sentence splitting
            return re.split(r'(?<=[.!?])\s+', text.strip())
    def _group_similar_sentences(self, sentences: List[str], embeddings: List[List[float]]) -> List[str]:
        chunks = []
        current_chunk = []
        current_embedding = None
        for sentence, embedding in zip(sentences, embeddings):
            if not current_chunk:
                current_chunk.append(sentence)
                current_embedding = embedding
            else:
                similarity = self._cosine_similarity(current_embedding, embedding)
                if similarity >= self.similarity_threshold:
                    current_chunk.append(sentence)
                    current_embedding = np.mean([current_embedding, embedding], axis=0)
                else:
                    chunks.append(" ".join(current_chunk))
                    current_chunk = [sentence]
                    current_embedding = embedding
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        return chunks
    def _cosine_similarity(self, vec1: List[float], vec2: List[float]) -> float:
        vec1 = np.array(vec1)
        vec2 = np.array(vec2)
        return np.dot(vec1, vec2) / (np.linalg.norm(vec1) * np.linalg.norm(vec2))
    def _calculate_semantic_score(self, text: str) -> float:
        return 1.0

class SentenceChunker:
    """Chunker that splits documents into sentences."""
    def chunk(self, text: str) -> List[str]:
        if _HAS_NLTK:
            return sent_tokenize(text)
        return re.split(r'(?<=[.!?])\s+', text.strip())

class SlidingWindowChunker:
    """Chunker that splits documents using a sliding window approach."""
    def chunk(self, text: str, window_size: int = 100, stride: int = 50) -> List[str]:
        words = text.split()
        chunks = []
        for i in range(0, len(words), stride):
            chunk = words[i:i+window_size]
            if chunk:
                chunks.append(' '.join(chunk))
            if i + window_size >= len(words):
                break
        return chunks

class RecursiveChunker:
    """Chunker that recursively splits documents by paragraphs, then sentences, then tokens."""
    def chunk(self, text: str, max_length: int = 512) -> List[str]:
        paragraphs = text.split('\n\n')
        chunks = []
        for para in paragraphs:
            if len(para.split()) <= max_length:
                chunks.append(para.strip())
            else:
                if _HAS_NLTK:
                    sentences = sent_tokenize(para)
                else:
                    sentences = re.split(r'(?<=[.!?])\s+', para.strip())
                current = []
                for sent in sentences:
                    if len(' '.join(current + [sent]).split()) <= max_length:
                        current.append(sent)
                    else:
                        if current:
                            chunks.append(' '.join(current))
                        current = [sent]
                if current:
                    chunks.append(' '.join(current))
        final_chunks = []
        for chunk in chunks:
            words = chunk.split()
            if len(words) > max_length:
                for i in range(0, len(words), max_length):
                    final_chunks.append(' '.join(words[i:i+max_length]))
            else:
                final_chunks.append(chunk)
        return [c for c in final_chunks if c.strip()]

class TokenChunker:
    """Chunker that splits text into chunks of N tokens using a HuggingFace tokenizer."""
    def __init__(self, tokenizer: Any, max_tokens: int = 512, stride: int = 256):
        self.tokenizer = tokenizer
        self.max_tokens = max_tokens
        self.stride = stride
    def chunk(self, text: str) -> List[str]:
        tokens = self.tokenizer.encode(text, add_special_tokens=False)
        chunks = []
        for i in range(0, len(tokens), self.stride):
            chunk_tokens = tokens[i:i+self.max_tokens]
            if not chunk_tokens:
                break
            chunk_text = self.tokenizer.decode(chunk_tokens)
            chunks.append(chunk_text)
            if i + self.max_tokens >= len(tokens):
                break
        return chunks

class OverlappingSentenceChunker:
    """Chunker that splits text into overlapping sentence windows."""
    def __init__(self, window_size: int = 5, stride: int = 2):
        self.window_size = window_size
        self.stride = stride
    def chunk(self, text: str) -> List[str]:
        try:
            from nltk.tokenize import sent_tokenize
            sentences = sent_tokenize(text)
        except ImportError:
            sentences = re.split(r'(?<=[.!?])\s+', text.strip())
        chunks = []
        for i in range(0, len(sentences), self.stride):
            chunk = sentences[i:i+self.window_size]
            if chunk:
                chunks.append(' '.join(chunk))
            if i + self.window_size >= len(sentences):
                break
        return chunks

class CodeChunker:
    """Chunker that splits code into logical blocks (functions, classes, etc.)."""
    def chunk(self, code: str) -> List[str]:
        pattern = re.compile(r'(^\s*def\s+|^\s*class\s+)', re.MULTILINE)
        indices = [m.start() for m in pattern.finditer(code)]
        indices.append(len(code))
        chunks = []
        for i in range(len(indices)-1):
            chunk = code[indices[i]:indices[i+1]].strip()
            if chunk:
                chunks.append(chunk)
        return chunks

class TableChunker:
    """Chunker that extracts and splits tables from text (e.g., markdown or CSV tables)."""
    def chunk(self, text: str) -> List[str]:
        tables = []
        md_table_pattern = re.compile(r'(\|.+\|\n)(\|[-: ]+\|\n)((\|.*\|\n)+)', re.MULTILINE)
        for match in md_table_pattern.finditer(text):
            tables.append(match.group())
        csv_lines = [line for line in text.splitlines() if ',' in line]
        if csv_lines:
            tables.append('\n'.join(csv_lines))
        return tables

class CharacterChunker:
    """Chunker that splits text into fixed-size character windows."""
    def __init__(self, window_size: int = 1000, stride: int = 1000):
        self.window_size = window_size
        self.stride = stride
    def chunk(self, text: str) -> List[str]:
        return [text[i:i+self.window_size] for i in range(0, len(text), self.stride) if text[i:i+self.window_size]]

class ParagraphChunker:
    """Chunker that splits text by paragraphs (double newlines or indentation)."""
    def chunk(self, text: str) -> List[str]:
        paras = re.split(r'(?:\n\s*\n|^\s+)', text, flags=re.MULTILINE)
        return [p.strip() for p in paras if p.strip()]

class LanguageSpecificChunker:
    """Chunker for Chinese/Japanese/Korean using language-specific tokenizers."""
    def __init__(self, language: str = 'zh'):
        self.language = language
        if language == 'zh':
            try:
                import jieba
                self.tokenizer = jieba
            except ImportError:
                raise ImportError('jieba is required for Chinese tokenization. Install with: pip install jieba')
    def chunk(self, text: str, window_size: int = 100, stride: int = 100) -> List[str]:
        if self.language == 'zh':
            tokens = list(self.tokenizer.cut(text))
            return [''.join(tokens[i:i+window_size]) for i in range(0, len(tokens), stride) if tokens[i:i+window_size]]
        return [text]

class HTMLXMLChunker:
    """Chunker that splits by HTML/XML tags or sections."""
    def __init__(self, tag: Optional[str] = None):
        self.tag = tag
    def chunk(self, html: str) -> List[str]:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError('BeautifulSoup4 is required for HTML/XML chunking. Install with: pip install beautifulsoup4')
        soup = BeautifulSoup(html, 'html.parser')
        if self.tag:
            return [str(e) for e in soup.find_all(self.tag)]
        return [str(e) for e in soup.body.find_all(recursive=False)] if soup.body else [str(e) for e in soup.find_all(recursive=False)]

class AudioVideoChunker:
    """Chunker for audio/video files using speech-to-text. User must provide a transcribe_fn callable."""
    def __init__(self, transcribe_fn: Optional[Callable[[str], str]] = None):
        self.transcribe_fn = transcribe_fn
    def chunk(self, file_path: str, window_size: int = 1000, stride: int = 1000) -> List[str]:
        if not self.transcribe_fn:
            raise ValueError('You must provide a transcribe_fn for audio/video chunking.')
        text = self.transcribe_fn(file_path)
        return [text[i:i+window_size] for i in range(0, len(text), stride) if text[i:i+window_size]]

class ImageChunker:
    """Chunker for images using OCR. User must provide an ocr_fn callable."""
    def __init__(self, ocr_fn: Optional[Callable[[str], str]] = None):
        self.ocr_fn = ocr_fn
    def chunk(self, image_path: str, window_size: int = 1000, stride: int = 1000) -> List[str]:
        if not self.ocr_fn:
            raise ValueError('You must provide an ocr_fn for image chunking.')
        text = self.ocr_fn(image_path)
        return [text[i:i+window_size] for i in range(0, len(text), stride) if text[i:i+window_size]]

class CustomChunker:
    """Chunker that allows user to pass a function or regex for custom chunking logic."""
    def __init__(self, chunk_fn: Optional[Callable[[str], List[str]]] = None, regex: Optional[str] = None):
        self.chunk_fn = chunk_fn
        self.regex = regex
    def chunk(self, text: str) -> List[str]:
        if self.chunk_fn:
            return self.chunk_fn(text)
        if self.regex:
            return re.split(self.regex, text)
        return [text]

class HybridChunker:
    """Chunker that combines multiple strategies (e.g., semantic + fixed size fallback)."""
    def __init__(self, chunkers: List[Any]):
        self.chunkers = chunkers
    def chunk(self, text: str) -> List[str]:
        for chunker in self.chunkers:
            chunks = chunker.chunk(text)
            if chunks and all(len(c) < 10000 for c in chunks):
                return chunks
        return [text]

class SectionHeadingChunker:
    """Chunker that splits by document headings (Markdown, HTML, DOCX, PDF headings)."""
    def __init__(self, heading_regex: Optional[str] = None):
        self.heading_regex = heading_regex or r'(^#+\s+.*$)'
    def chunk(self, text: str) -> List[str]:
        splits = re.split(self.heading_regex, text, flags=re.MULTILINE)
        chunks = []
        for i in range(1, len(splits), 2):
            heading = splits[i].strip()
            content = splits[i+1].strip() if i+1 < len(splits) else ''
            chunks.append(f'{heading}\n{content}')
        return chunks if chunks else [text]

class OverlappingParagraphChunker:
    """Chunker that splits text into overlapping paragraph windows."""
    def __init__(self, window_size: int = 3, stride: int = 1):
        self.window_size = window_size
        self.stride = stride
    def chunk(self, text: str) -> List[str]:
        paras = [p.strip() for p in re.split(r'(?:\n\s*\n|^\s+)', text, flags=re.MULTILINE) if p.strip()]
        chunks = []
        for i in range(0, len(paras), self.stride):
            chunk = paras[i:i+self.window_size]
            if chunk:
                chunks.append('\n\n'.join(chunk))
            if i + self.window_size >= len(paras):
                break
        return chunks 

@dataclass
class DocumentChunk:
    """Represents a processed document chunk."""
    text: str
    metadata: Dict[str, Any]
    chunk_id: str
    parent_id: Optional[str]
    semantic_score: Optional[float] = None
    embedding: Optional[List[float]] = None

class ChunkingStrategy(Enum):
    """Different document chunking strategies."""
    FIXED_SIZE = "fixed_size"
    SEMANTIC = "semantic"
    RECURSIVE = "recursive"
    SLIDING_WINDOW = "sliding_window"

class MetadataExtractor:
    """Extracts and enriches document metadata."""

    def __init__(self, nlp_model: Optional[str] = "en_core_web_sm"):
        if SPACY_AVAILABLE and nlp_model:
            self.nlp = spacy.load(nlp_model)
        else:
            self.nlp = None

    def extract_metadata(self, text: str) -> Dict[str, Any]:
        """
        Extract metadata from text using NLP.
        
        Args:
            text: Input text
            
        Returns:
            Dictionary of extracted metadata
        """
        if not self.nlp:
            return {}

        doc = self.nlp(text)
        
        # Extract entities
        entities = {
            ent.label_: [e.text for e in doc.ents if e.label_ == ent.label_]
            for ent in doc.ents
        }
        
        # Extract key phrases (noun chunks)
        key_phrases = [chunk.text for chunk in doc.noun_chunks]
        
        # Extract document statistics
        stats = {
            "word_count": len(doc),
            "sentence_count": len(list(doc.sents)),
            "avg_word_length": np.mean([len(token.text) for token in doc]),
            "unique_words": len(set(token.text.lower() for token in doc))
        }
        
        return {
            "entities": entities,
            "key_phrases": key_phrases,
            "statistics": stats
        }

class SpreadsheetChunker:
    """Chunker for spreadsheets (Excel/CSV). Splits by rows, columns, or sheets."""
    def __init__(self, mode: str = 'row', sheet_name: str = None):
        """
        mode: 'row', 'column', or 'sheet'
        sheet_name: for Excel, specify a sheet to chunk
        """
        self.mode = mode
        self.sheet_name = sheet_name
    def chunk(self, file_path: str) -> list:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError('pandas is required for SpreadsheetChunker. Install with: pip install pandas openpyxl')
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path, sheet_name=self.sheet_name)
        if self.mode == 'row':
            return [row.to_json() for _, row in df.iterrows()]
        elif self.mode == 'column':
            return [df[col].to_json() for col in df.columns]
        elif self.mode == 'sheet':
            if hasattr(df, 'items'):
                # Multiple sheets
                return [sheet_df.to_json() for _, sheet_df in df.items()]
            else:
                return [df.to_json()]
        else:
            raise ValueError('mode must be one of: row, column, sheet')

class PresentationChunker:
    """Chunker for presentations (PowerPoint). Splits by slide."""
    def chunk(self, file_path: str) -> list:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError('python-pptx is required for PresentationChunker. Install with: pip install python-pptx')
        prs = Presentation(file_path)
        slides = []
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            slides.append("\n".join(text))
        return slides

class AdaptiveHybridChunker:
    """Hybrid chunker that selects chunking strategy based on content type (text, table, code, etc.)."""
    def __init__(self, chunker_map: dict):
        """
        chunker_map: dict mapping content type (e.g., 'text', 'table', 'code', 'spreadsheet', 'presentation') to chunker instance
        """
        self.chunker_map = chunker_map
    def chunk(self, content: Any, content_type: str = 'text', **kwargs) -> list:
        if content_type not in self.chunker_map:
            raise ValueError(f'No chunker registered for content type: {content_type}')
        chunker = self.chunker_map[content_type]
        return chunker.chunk(content, **kwargs)
