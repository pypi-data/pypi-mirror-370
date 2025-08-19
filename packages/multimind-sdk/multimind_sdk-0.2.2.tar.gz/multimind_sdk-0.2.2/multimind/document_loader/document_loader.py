"""
Enhanced document loading with support for multiple formats and sources.
"""

from typing import List, Dict, Any, Optional, Union, Protocol, runtime_checkable, Tuple, Callable
from pathlib import Path
import asyncio
import aiohttp
from dataclasses import dataclass
from enum import Enum
import json
import logging
try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
try:
    import docx
except ImportError:
    docx = None
try:
    import pandas as pd
except ImportError:
    pd = None
try:
    from unstructured.partition.auto import partition
except ImportError:
    partition = None
from ..models.base import BaseLLM
import os

@dataclass
class DocumentMetadata:
    """Metadata for loaded documents."""
    source: str
    format: str
    created_at: Optional[str] = None
    modified_at: Optional[str] = None
    author: Optional[str] = None
    title: Optional[str] = None
    page_number: Optional[int] = None
    custom_metadata: Optional[Dict[str, Any]] = None

@dataclass
class LoadedDocument:
    """Represents a loaded document with content and metadata."""
    content: str
    metadata: DocumentMetadata
    raw_content: Optional[Any] = None  # Original format content

class DocumentFormat(Enum):
    """Supported document formats."""
    PDF = "pdf"
    DOCX = "docx"
    TXT = "txt"
    HTML = "html"
    JSON = "json"
    CSV = "csv"
    MARKDOWN = "md"
    UNSTRUCTURED = "unstructured"

class DocumentSource(Enum):
    """Supported document sources."""
    LOCAL = "local"
    URL = "url"
    DATABASE = "database"
    API = "api"
    STREAM = "stream"

@runtime_checkable
class DocumentConnector(Protocol):
    """Protocol for document connectors."""
    async def connect(self) -> None:
        """Establish connection to the document source."""
        ...
    
    async def disconnect(self) -> None:
        """Close connection to the document source."""
        ...
    
    async def fetch_documents(self, **kwargs) -> List[LoadedDocument]:
        """Fetch documents from the source."""
        ...

class BaseDocumentLoader:
    """Base class for document loaders."""

    def __init__(self, **kwargs):
        self.kwargs = kwargs
        self._semaphore = asyncio.Semaphore(kwargs.get('max_concurrent_operations', 10))

    async def _execute_with_semaphore(self, coro):
        """Execute coroutine with semaphore for rate limiting."""
        async with self._semaphore:
            return await coro

    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        """Load a single document. Must be implemented in subclass."""
        raise NotImplementedError("load_document must be implemented in a subclass of BaseDocumentLoader.")

    async def load_documents(self, sources: List[str], **kwargs) -> List[LoadedDocument]:
        """Load multiple documents in parallel."""
        tasks = [self.load_document(source, **kwargs) for source in sources]
        return await asyncio.gather(*tasks)

class LocalDocumentLoader(BaseDocumentLoader):
    """Loader for local documents."""

    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        """Load a document from local filesystem."""
        try:
            path = Path(source)
            if not path.exists():
                raise FileNotFoundError(f"Document not found: {source}")

            format = DocumentFormat(path.suffix[1:].lower())
            metadata = DocumentMetadata(
                source=str(path),
                format=format.value,
                created_at=str(path.stat().st_ctime),
                modified_at=str(path.stat().st_mtime)
            )

            if format == DocumentFormat.PDF:
                content, raw = await self._load_pdf(path)
            elif format == DocumentFormat.DOCX:
                content, raw = await self._load_docx(path)
            elif format == DocumentFormat.TXT:
                content, raw = await self._load_txt(path)
            elif format == DocumentFormat.HTML:
                content, raw = await self._load_html(path)
            elif format == DocumentFormat.JSON:
                content, raw = await self._load_json(path)
            elif format == DocumentFormat.CSV:
                content, raw = await self._load_csv(path)
            elif format == DocumentFormat.MARKDOWN:
                content, raw = await self._load_markdown(path)
            else:
                content, raw = await self._load_unstructured(path)

            return LoadedDocument(
                content=content,
                metadata=metadata,
                raw_content=raw
            )

        except Exception as e:
            logging.error(f"Error loading document {source}: {str(e)}")
            raise

    async def _load_pdf(self, path: Path) -> Tuple[str, Any]:
        """Load PDF document."""
        with open(path, 'rb') as f:
            pdf = PyPDF2.PdfReader(f)
            content = []
            raw = pdf
            for page in pdf.pages:
                content.append(page.extract_text())
            return "\n".join(content), raw

    async def _load_docx(self, path: Path) -> Tuple[str, Any]:
        """Load DOCX document."""
        doc = docx.Document(path)
        content = []
        raw = doc
        for para in doc.paragraphs:
            content.append(para.text)
        return "\n".join(content), raw

    async def _load_txt(self, path: Path) -> Tuple[str, Any]:
        """Load text document."""
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
            return content, content

    async def _load_html(self, path: Path) -> Tuple[str, Any]:
        """Load HTML document."""
        with open(path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            content = soup.get_text(separator='\n')
            return content, soup

    async def _load_json(self, path: Path) -> Tuple[str, Any]:
        """Load JSON document."""
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            content = json.dumps(data, indent=2)
            return content, data

    async def _load_csv(self, path: Path) -> Tuple[str, Any]:
        """Load CSV document."""
        df = pd.read_csv(path)
        content = df.to_string()
        return content, df

    async def _load_markdown(self, path: Path) -> Tuple[str, Any]:
        """Load Markdown document."""
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
            return content, content

    async def _load_unstructured(self, path: Path) -> Tuple[str, Any]:
        """Load document using unstructured."""
        elements = partition(str(path))
        content = "\n".join([str(el) for el in elements])
        return content, elements

class WebDocumentLoader(BaseDocumentLoader):
    """Loader for web documents."""

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.session = None

    async def _ensure_session(self):
        """Ensure aiohttp session exists."""
        if self.session is None:
            self.session = aiohttp.ClientSession()

    async def load_document(self, url: str, **kwargs) -> LoadedDocument:
        """Load a document from URL."""
        try:
            await self._ensure_session()
            async with self.session.get(url) as response:
                if response.status != 200:
                    raise ValueError(f"Failed to fetch document: {url}")

                content_type = response.headers.get('content-type', '')
                if 'application/pdf' in content_type:
                    content, raw = await self._load_pdf_from_url(response)
                elif 'application/json' in content_type:
                    content, raw = await self._load_json_from_url(response)
                elif 'text/html' in content_type:
                    content, raw = await self._load_html_from_url(response)
                else:
                    content, raw = await self._load_text_from_url(response)

                metadata = DocumentMetadata(
                    source=url,
                    format=content_type.split(';')[0],
                    modified_at=response.headers.get('last-modified')
                )

                return LoadedDocument(
                    content=content,
                    metadata=metadata,
                    raw_content=raw
                )

        except Exception as e:
            logging.error(f"Error loading document from {url}: {str(e)}")
            raise

    async def _load_pdf_from_url(self, response: aiohttp.ClientResponse) -> Tuple[str, Any]:
        """Load PDF from URL."""
        content = await response.read()
        pdf = PyPDF2.PdfReader(io.BytesIO(content))
        text_content = []
        for page in pdf.pages:
            text_content.append(page.extract_text())
        return "\n".join(text_content), pdf

    async def _load_json_from_url(self, response: aiohttp.ClientResponse) -> Tuple[str, Any]:
        """Load JSON from URL."""
        data = await response.json()
        return json.dumps(data, indent=2), data

    async def _load_html_from_url(self, response: aiohttp.ClientResponse) -> Tuple[str, Any]:
        """Load HTML from URL."""
        html = await response.text()
        soup = BeautifulSoup(html, 'html.parser')
        return soup.get_text(separator='\n'), soup

    async def _load_text_from_url(self, response: aiohttp.ClientResponse) -> Tuple[str, Any]:
        """Load text from URL."""
        content = await response.text()
        return content, content

    async def __aenter__(self):
        """Context manager entry."""
        await self._ensure_session()
        return self

    async def __aexit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit."""
        if self.session:
            await self.session.close()
            self.session = None

class DatabaseDocumentLoader(BaseDocumentLoader):
    """Loader for database documents."""

    def __init__(self, connector: DocumentConnector, **kwargs):
        super().__init__(**kwargs)
        self.connector = connector

    async def load_documents(self, **kwargs) -> List[LoadedDocument]:
        """Load documents from database."""
        try:
            await self.connector.connect()
            return await self.connector.fetch_documents(**kwargs)
        finally:
            await self.connector.disconnect()

class StreamDocumentLoader(BaseDocumentLoader):
    """Loader for streaming documents."""

    def __init__(self, stream_connector: DocumentConnector, **kwargs):
        super().__init__(**kwargs)
        self.connector = stream_connector
        self._stream_task = None

    async def start_streaming(self, callback: Callable[[LoadedDocument], None], **kwargs):
        """Start streaming documents."""
        try:
            await self.connector.connect()
            self._stream_task = asyncio.create_task(
                self._stream_documents(callback, **kwargs)
            )
        except Exception as e:
            logging.error(f"Error starting stream: {str(e)}")
            raise

    async def stop_streaming(self):
        """Stop streaming documents."""
        if self._stream_task:
            self._stream_task.cancel()
            try:
                await self._stream_task
            except asyncio.CancelledError:
                pass
            finally:
                await self.connector.disconnect()

    async def _stream_documents(
        self,
        callback: Callable[[LoadedDocument], None],
        **kwargs
    ):
        """Stream documents to callback."""
        try:
            async for doc in self.connector.stream_documents(**kwargs):
                await callback(doc)
        except asyncio.CancelledError:
            raise
        except Exception as e:
            logging.error(f"Error streaming documents: {str(e)}")
            raise

class DocumentLoaderFactory:
    """Factory for creating document loaders."""

    @staticmethod
    def create_loader(
        source_type: DocumentSource,
        **kwargs
    ) -> BaseDocumentLoader:
        """Create appropriate document loader."""
        if source_type == DocumentSource.LOCAL:
            return LocalDocumentLoader(**kwargs)
        elif source_type == DocumentSource.URL:
            return WebDocumentLoader(**kwargs)
        elif source_type == DocumentSource.DATABASE:
            return DatabaseDocumentLoader(kwargs.pop('connector'), **kwargs)
        elif source_type == DocumentSource.STREAM:
            return StreamDocumentLoader(kwargs.pop('connector'), **kwargs)
        else:
            raise ValueError(f"Unsupported source type: {source_type}")

class WebsiteDocumentLoader:
    """Loader for ingesting documents from websites (HTML/webpages)."""
    async def load(self, url: str) -> Tuple[str, str]:
        """Fetch and extract main text content from a webpage."""
        try:
            import requests
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("Please install 'requests' and 'beautifulsoup4' to use WebsiteDocumentLoader.")
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        # Try to extract main content
        texts = [t for t in soup.stripped_strings]
        content = '\n'.join(texts)
        return content, response.text

class EmailDocumentLoader:
    """Loader for parsing and ingesting email files (EML, MSG, etc.)."""
    async def load(self, file_path: str) -> Tuple[str, str]:
        """Parse an email file and extract the main text content."""
        import email
        from email import policy
        from email.parser import BytesParser
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Email file not found: {file_path}")
        with open(file_path, 'rb') as f:
            msg = BytesParser(policy=policy.default).parse(f)
        # Extract text/plain part
        text = ""
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == 'text/plain':
                    text += part.get_content()
        else:
            text = msg.get_content()
        return text.strip(), str(msg)

class SpreadsheetDocumentLoader(BaseDocumentLoader):
    """Loader for spreadsheet documents (Excel/CSV)."""
    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError('pandas is required for SpreadsheetDocumentLoader. Install with: pip install pandas openpyxl')
        path = Path(source)
        if path.suffix.lower() == '.csv':
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
        content = df.to_string()
        metadata = DocumentMetadata(source=str(path), format=path.suffix[1:].lower())
        return LoadedDocument(content=content, metadata=metadata, raw_content=df)

class PresentationDocumentLoader(BaseDocumentLoader):
    """Loader for presentation documents (PowerPoint)."""
    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError('python-pptx is required for PresentationDocumentLoader. Install with: pip install python-pptx')
        path = Path(source)
        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            slides.append("\n".join(text))
        content = "\n---\n".join(slides)
        metadata = DocumentMetadata(source=str(path), format=path.suffix[1:].lower())
        return LoadedDocument(content=content, metadata=metadata, raw_content=prs)

class ImageDocumentLoader(BaseDocumentLoader):
    """Loader for image files (extracts text via OCR)."""
    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            raise ImportError('Pillow and pytesseract are required for ImageDocumentLoader. Install with: pip install pillow pytesseract')
        path = Path(source)
        image = Image.open(path)
        content = pytesseract.image_to_string(image)
        metadata = DocumentMetadata(source=str(path), format=path.suffix[1:].lower())
        return LoadedDocument(content=content, metadata=metadata, raw_content=image)

class AudioDocumentLoader(BaseDocumentLoader):
    """Loader for audio files (extracts text via speech-to-text)."""
    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        try:
            import librosa
        except ImportError:
            raise ImportError('librosa is required for AudioDocumentLoader. Install with: pip install librosa')
        # User must provide a transcribe_fn for actual speech-to-text
        transcribe_fn = kwargs.get('transcribe_fn')
        if not transcribe_fn:
            raise ValueError('You must provide a transcribe_fn for audio transcription.')
        path = Path(source)
        audio, sr = librosa.load(path, sr=None)
        content = transcribe_fn(audio, sr)
        metadata = DocumentMetadata(source=str(path), format=path.suffix[1:].lower())
        return LoadedDocument(content=content, metadata=metadata, raw_content=audio)

class VideoDocumentLoader(BaseDocumentLoader):
    """Loader for video files (extracts text via video-to-text or speech-to-text)."""
    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        try:
            import moviepy.editor as mp
        except ImportError:
            raise ImportError('moviepy is required for VideoDocumentLoader. Install with: pip install moviepy')
        # User must provide a transcribe_fn for actual video/audio transcription
        transcribe_fn = kwargs.get('transcribe_fn')
        if not transcribe_fn:
            raise ValueError('You must provide a transcribe_fn for video transcription.')
        path = Path(source)
        video = mp.VideoFileClip(str(path))
        audio = video.audio
        audio_path = str(path) + '.temp_audio.wav'
        audio.write_audiofile(audio_path)
        import librosa
        audio_data, sr = librosa.load(audio_path, sr=None)
        content = transcribe_fn(audio_data, sr)
        os.remove(audio_path)
        metadata = DocumentMetadata(source=str(path), format=path.suffix[1:].lower())
        return LoadedDocument(content=content, metadata=metadata, raw_content=video)

class DefaultFileLoader(BaseDocumentLoader):
    """Default file loader that loads text files from disk."""
    async def load_document(self, source: str, **kwargs) -> LoadedDocument:
        if not os.path.isfile(source):
            raise FileNotFoundError(f"File not found: {source}")
        with open(source, "r", encoding="utf-8") as f:
            text = f.read()
        return LoadedDocument(text=text, metadata={"source": source}) 