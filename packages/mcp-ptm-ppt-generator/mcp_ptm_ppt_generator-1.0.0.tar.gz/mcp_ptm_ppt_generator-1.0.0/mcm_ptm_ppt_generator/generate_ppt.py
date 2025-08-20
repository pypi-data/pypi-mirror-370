#!/usr/bin/env python3

from fastmcp import FastMCP
from pydantic import Field, BaseModel
from pptx import Presentation
import os

mcp = FastMCP("ptms made easier")

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "template.pptx")
OUTPUT_FOLDER = os.path.join(SCRIPT_DIR, "generated_ppts")
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


class Student(BaseModel):
    Tutor: str = Field(..., description="Tutor's name:")
    Student: str = Field(..., description="Student's name:")
    Subjects: str = Field(..., description="Student's subject:")
    ParentRequirement: str = Field(..., description="Parent Requirement:")
    ReportingPeriod: str = Field(..., description="Reoprting period:")
    NoOfSessions: str = Field(..., description="No of sessions:")
    Target: str = Field(..., description="IXL Target score")
    Numbers: str = Field(..., description="IXL Numbers and Operations Score")
    Algebra: str = Field(..., description="IXL Algebra and Algebraic Thinking Score")
    Fractions: str = Field(..., description="IXL Fractions Score")
    Geometry: str = Field(..., description="IXL Goemetry Score")
    Measurement: str = Field(..., description="IXL Measurement Score")
    Data: str = Field(..., description="IXL Data & Probability Score")
    Overall: str = Field(..., description="Overall IXL Math Level")
    #  add the fields for the third slide here to talk about the rec skills from IXL
    IXLAreaOfImprovement1: str = Field(..., description="Suggested IXL strand to work on:")
    IXLAreaOfImprovement2: str = Field(..., description="IXL Suggested strand 2 to work on:")
    AreaOfImprovement1SuggestedSkill1: str = Field(..., description="Skill 1 to work on to improve area of improvement 1")
    AreaOfImprovement1SuggestedSkill2: str = Field(..., description="Skill 2 to work on to improve area of improvement 1")
    AreaOfImprovement2SuggestedSkill1: str = Field(..., description="Skill 1 to work on to improve area of improvement 2")
    AreaOfImprovement2SuggestedSkill2: str = Field(..., description="Skill 2 to work on to improve area of improvement 2")
    # topics covered stuff
    Topic1: str = Field(..., description="Topic 1 covered this month")
    T1Status: str = Field(..., description="Status of the topic 1")
    Topic2: str = Field(..., description="Topic 2 covered this month")
    T2Status: str = Field(..., description="Status of the topic 2")
    MTest: str = Field(..., description="Monthly Test Score (out of 25)")
    LGap: str = Field(..., description="Learning gap identified")
    APlan: str = Field(..., description="Action Plan for the learning Gap")
    StudentStepsNeeded: str = Field(..., description="Steps needed from student")
    Task1: str = Field(..., description="Next task planned")
    Task1Sess: str = Field(..., description="Number of sessions needed for task 1")
    Task2: str = Field(..., description="Next task2 planned")
    Task2Sess: str = Field(..., description="Number of sessions needed for task 2")
    Notes: str = Field(..., description="Notes for upcoming tasks planned")





def replace_text(shape, data: dict):
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in data.items():
                    placeholder = f"{{{{{key}}}}}"
                    if placeholder in run.text:
                        if key == "Subjects":
                            value = ", ".join([s.strip() for s in value.split(",")])
                        run.text = run.text.replace(placeholder, str(value))
    if shape.shape_type == 6:
        for subshape in shape.shapes:
            replace_text(subshape, data)
    


    if shape.shape_type == 19:
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in data.items():
                            placeholder= f"{{{{{key}}}}}"
                            if placeholder in run.text:
                                if key == "Subjects":
                                    value = ", ".join([s.strip() for s in value.split(",")])
                                run.text = run.text.replace(placeholder, str(value))

def add_stuff_to_ppt(prs: Presentation, data: Student):
    d = data.dict()
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text(shape, d)



@mcp.tool
def generate_ppt(student: Student):
    """Genrate a ppt for your lovely students based on the template"""
    prs = Presentation(TEMPLATE_PATH)
    add_stuff_to_ppt(prs, student)
    output_path = os.path.join(OUTPUT_FOLDER, f"{student.Student}.pptx")
    prs.save(output_path)
    return f"PPT generated for {student.Student} at {output_path}"



def main():
    mcp.run()

if __name__=="__main__":
    main()

