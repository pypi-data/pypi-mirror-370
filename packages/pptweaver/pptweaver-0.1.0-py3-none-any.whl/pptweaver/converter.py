import asyncio
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
import logging
import os
import re
from bs4 import BeautifulSoup
from pptx.enum.shapes import MSO_SHAPE

# Import processors
from .elements.rect_processor import RectProcessor
from .elements.text_processor import TextProcessor
from .elements.line_processor import LineProcessor
from .elements.polygon_processor import PolygonProcessor
from .elements.oval_processor import OvalProcessor
from .elements.path_processor import PathProcessor
from .elements.foreign_object_processor import ForeignObjectProcessor
from .elements.image_processor import ImageProcessor

# --- Constants ---
PPTX_WIDTH = Emu(9144000)
PPTX_HEIGHT = Emu(5143500)
EMU_PER_PIXEL = 914400 / 96

class Converter:
    def __init__(self, input_file: str, output_file: str):
        self.input_file_path = f"file://{os.path.abspath(input_file)}"
        self.output_file = output_file
        self.prs = Presentation()
        self.prs.slide_width = PPTX_WIDTH
        self.prs.slide_height = PPTX_HEIGHT
        self.page = None
        self.scale_x = 1.0
        self.scale_y = 1.0
        self._register_processors()

    def _register_processors(self):
        self._processors = {
            # SVG Shapes
            "rect": RectProcessor,
            "line": LineProcessor,
            "polygon": PolygonProcessor,
            "polyline": PolygonProcessor,
            "circle": OvalProcessor,
            "ellipse": OvalProcessor,
            "path": PathProcessor,
            "foreignobject": ForeignObjectProcessor,
            "image": ImageProcessor,
            
            # Text Elements
            "text": TextProcessor,
            # The following HTML tags are handled by the ForeignObjectProcessor's rich text parsing
            # "p": TextProcessor,
            # "li": TextProcessor,
            # "h3": TextProcessor,
            # "div": TextProcessor,
        }

    # --- Utility methods now part of the class ---

    def px_to_emu(self, px):
        return int(px * EMU_PER_PIXEL)

    def px_to_pt(self, px):
        return Pt(px * 0.75)

    def parse_color(self, color_string):
        if not color_string or color_string in ['none', 'transparent']:
            return None, None
        
        match = re.match(r'rgba\((\d+),\s*(\d+),\s*(\d+),\s*([\d.]+)\)', color_string)
        if match:
            color = RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3)))
            alpha = float(match.group(4))
            return color, alpha

        match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_string)
        if match:
            return RGBColor(int(match.group(1)), int(match.group(2)), int(match.group(3))), 1.0
        
        if color_string.startswith('#') and len(color_string) == 7:
            return RGBColor.from_string(color_string[1:]), 1.0
        
        if color_string.startswith('#') and len(color_string) == 4:
            r, g, b = color_string[1], color_string[2], color_string[3]
            return RGBColor.from_string(f'{r}{r}{g}{g}{b}{b}'), 1.0
        
        logging.warning(f"Unsupported color format: {color_string}. Defaulting to black.")
        return RGBColor(0, 0, 0), 1.0

    async def convert(self):
        # Read the file content once
        with open(self.input_file_path.replace('file://', ''), 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        # Use a more specific selector to find only direct children of <html>
        slides_content = soup.select('html > body')

        if not slides_content:
            logging.error("No <body> tags found in the input file.")
            return

        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            self.page = await browser.new_page()

            logging.info(f"Found {len(slides_content)} slides to convert.")

            for i, body_tag in enumerate(slides_content):
                logging.info(f"Processing slide {i+1}...")
                # Set the content of the page to the current body tag
                await self.page.set_content(str(body_tag))
                slide_handle = await self.page.query_selector('body')
                
                if slide_handle:
                    await self._process_slide(slide_handle)

            await browser.close()
        
        logging.info(f"Saving presentation to {self.output_file}")
        self.prs.save(self.output_file)
        
    async def _process_slide(self, slide_handle):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Set slide background color
        bg_color_str = await slide_handle.evaluate('(element) => window.getComputedStyle(element).backgroundColor')
        bg_color, _ = self.parse_color(bg_color_str)
        if bg_color:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg_color

        slide_box = await slide_handle.bounding_box()
        if not slide_box or slide_box['width'] == 0 or slide_box['height'] == 0:
            logging.warning(f"Slide has zero dimensions. Skipping.")
            return
            
        self.scale_x = PPTX_WIDTH / self.px_to_emu(slide_box['width'])
        self.scale_y = PPTX_HEIGHT / self.px_to_emu(slide_box['height'])

        elements_data = await slide_handle.evaluate("""
            (slideElement) => {
                const elements = [];
                const svg = slideElement.querySelector('svg');
                if (!svg) return elements;

                const svgRect = svg.getBoundingClientRect();

                function processNode(node, parentTransform) {
                    if (node.nodeType !== Node.ELEMENT_NODE) return;

                    const tagName = node.tagName.toLowerCase();
                    const computedStyle = window.getComputedStyle(node);
                    
                    // Skip hidden elements
                    if (computedStyle.display === 'none' || computedStyle.visibility === 'hidden') {
                        return;
                    }

                    const rect = node.getBoundingClientRect();
                    
                    // Calculate position relative to the SVG container
                    const relRect = {
                        x: rect.x - svgRect.x,
                        y: rect.y - svgRect.y,
                        width: rect.width,
                        height: rect.height
                    };

                    const data = {
                        tagName: tagName,
                        rect: relRect,
                        attributes: {},
                        style: {
                            fill: computedStyle.fill,
                            stroke: computedStyle.stroke,
                            strokeWidth: computedStyle.strokeWidth,
                            fontSize: computedStyle.fontSize,
                            fontWeight: computedStyle.fontWeight,
                            fontFamily: computedStyle.fontFamily,
                            color: computedStyle.color,
                            opacity: computedStyle.opacity,
                            textAlign: computedStyle.textAlign,
                            markerEnd: computedStyle.markerEnd,
                        },
                        text: node.textContent ? node.textContent.trim() : ""
                    };

                    // Rich Text Extraction for foreignObject content
                    if (tagName === 'foreignobject') {
                        data.text_runs = [];
                        function extractRuns(element) {
                            element.childNodes.forEach(child => {
                                if (child.nodeType === Node.TEXT_NODE && child.textContent.trim()) {
                                    const style = window.getComputedStyle(child.parentNode);
                                    data.text_runs.push({
                                        text: child.textContent.trim(),
                                        tagName: child.parentNode.tagName.toLowerCase(),
                                        style: {
                                            fontWeight: style.fontWeight,
                                            fontFamily: style.fontFamily,
                                            color: style.color,
                                        }
                                    });
                                } else if (child.nodeType === Node.ELEMENT_NODE) {
                                    // Handle cases like <li><b>Text</b></li>
                                    if (child.childNodes.length > 0) {
                                        extractRuns(child);
                                    } else {
                                        const style = window.getComputedStyle(child);
                                        data.text_runs.push({
                                            text: child.textContent.trim(),
                                            tagName: child.tagName.toLowerCase(),
                                            style: {
                                                fontWeight: style.fontWeight,
                                                fontFamily: style.fontFamily,
                                                color: style.color,
                                            }
                                        });
                                    }
                                }
                            });
                        }
                        extractRuns(node);
                    }

                    // For image tags, capture the href
                    if (tagName === 'image') {
                        data.attributes['href'] = node.getAttribute('href');
                    }

                    // Special handling for <path> elements to linearize them
                    if (tagName === 'path') {
                        const pathLength = node.getTotalLength();
                        const points = [];
                        // Sample the path at intervals. A smaller step gives more detail.
                        const step = Math.max(2, pathLength / 100); 
                        for (let i = 0; i <= pathLength; i += step) {
                            const pt = node.getPointAtLength(i);
                            points.push(`${pt.x - svgRect.x},${pt.y - svgRect.y}`);
                        }
                        data.attributes['linearized_points'] = points.join(' ');
                    }

                    // Get attributes from the element
                    for (const attr of node.attributes) {
                        data.attributes[attr.name] = attr.value;
                    }

                    elements.push(data);

                    // Recursively process children
                    node.childNodes.forEach(child => processNode(child, "")); // Simplified transform for now
                }

                processNode(svg, "");
                return elements;
            }
        """)

        self._create_shapes_from_data(elements_data, slide)

    def _create_shapes_from_data(self, elements_data, slide):
        for data in elements_data:
            tag_name = data.get('tagName')
            processor_class = self._processors.get(tag_name)
            
            if processor_class:
                processor = processor_class(data, slide, self)
                try:
                    processor.process()
                except Exception as e:
                    logging.error(f"Error processing element <{tag_name}>: {e}", exc_info=True)
            elif tag_name and tag_name not in ['svg', 'g', 'defs', 'marker', 'polygon', 'ul', 'image', 'p', 'li', 'h3', 'div', 'b', 'span']:
                # Log only if it's a tag we might want to support, ignore structural/container tags
                logging.warning(f"No processor found for element <{tag_name}>. Skipping.")
