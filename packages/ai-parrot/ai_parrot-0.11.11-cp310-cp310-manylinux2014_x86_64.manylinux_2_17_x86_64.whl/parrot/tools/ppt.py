"""
PowerPoint Tool migrated to use AbstractDocumentTool framework.
"""
import re
from pathlib import Path
from typing import Dict, List, Optional, Any, Union
import io
import traceback
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from jinja2 import Environment, FileSystemLoader
from pydantic import Field, field_validator
import markdown
from bs4 import BeautifulSoup, NavigableString
from .document import (
    AbstractDocumentTool,
    DocumentGenerationArgs
)

class PowerPointArgs(DocumentGenerationArgs):
    """Arguments schema for PowerPoint presentation generation."""

    template_name: Optional[str] = Field(
        None,
        description="Name of the HTML template (e.g., 'presentation.html') to render before conversion"
    )
    template_vars: Optional[Dict[str, Any]] = Field(
        None,
        description="Variables to pass to the HTML template (e.g., title, author, date)"
    )
    pptx_template: Optional[str] = Field(
        None,
        description="Path to a PowerPoint template file (.pptx or .potx) to use as base"
    )
    slide_layout: int = Field(
        1,
        description="Default slide layout index (0=Title Slide, 1=Title and Content, etc.)",
        ge=0,
        le=15
    )
    title_styles: Optional[Dict[str, Any]] = Field(
        None,
        description="Styles to apply to slide titles (font_name, font_size, bold, italic, font_color, alignment)"
    )
    content_styles: Optional[Dict[str, Any]] = Field(
        None,
        description="Styles to apply to slide content (font_name, font_size, bold, italic, font_color, alignment)"
    )
    max_slides: int = Field(
        50,
        description="Maximum number of slides to generate",
        ge=1,
        le=100
    )
    split_by_headings: bool = Field(
        True,
        description="Whether to split content into slides based on headings (H1, H2, etc.)"
    )

    @field_validator('template_name')
    @classmethod
    def validate_template_name(cls, v):
        if v and not v.endswith('.html'):
            v = f"{v}.html"
        return v


class PowerPointTool(AbstractDocumentTool):
    """
    PowerPoint Presentation Generator Tool.

    This tool converts text content (including Markdown and HTML) into professionally
    formatted PowerPoint presentations. It automatically splits content into slides
    based on headings and supports custom templates, styling, and layout options.

    Features:
    - Automatic slide splitting based on headings (H1, H2, H3, etc.)
    - Markdown to PowerPoint conversion with proper formatting
    - HTML to PowerPoint conversion support
    - Custom PowerPoint template support
    - Jinja2 HTML template processing
    - Configurable slide layouts and styling
    - Table, list, and content formatting
    - Professional presentation generation

    Slide Splitting Logic:
    - H1 (# Title) → Title slide (layout 0)
    - H2 (## Section) → Content slide (layout 1)
    - H3 (### Subsection) → Content slide (layout 1)
    - Content between headings → Added to the slide
    """

    name = "powerpoint_generator"
    description = (
        "Generate PowerPoint presentations from text, Markdown, or HTML content. "
        "Automatically splits content into slides based on headings. "
        "Supports custom templates, styling, and professional presentation formatting."
    )
    args_schema = PowerPointArgs

    # Document type configuration
    document_type = "presentation"
    default_extension = "pptx"
    supported_extensions = [".pptx", ".potx"]

    def __init__(
        self,
        templates_dir: Optional[Path] = None,
        output_dir: Optional[Union[str, Path]] = None,
        default_html_template: Optional[str] = None,
        **kwargs
    ):
        """
        Initialize the PowerPoint Tool.

        Args:
            templates_dir: Directory containing HTML and PowerPoint templates
            output_dir: Directory where generated presentations will be saved
            default_html_template: Default HTML template for content processing
            **kwargs: Additional arguments for AbstractDocumentTool
        """
        # Set up output directory before calling super().__init__
        if output_dir:
            kwargs['output_dir'] = Path(output_dir)

        super().__init__(templates_dir=templates_dir, **kwargs)

        self.default_html_template = default_html_template

        # Initialize Jinja2 environment for HTML templates
        if self.templates_dir:
            self.html_env = Environment(
                loader=FileSystemLoader(str(self.templates_dir)),
                autoescape=True
            )
        else:
            self.html_env = None

    def _render_html_template(self, content: str, template_name: Optional[str], template_vars: Optional[Dict[str, Any]]) -> str:
        """Render content through Jinja2 HTML template if provided."""
        if not template_name or not self.html_env:
            return content

        try:
            template = self.html_env.get_template(template_name)
            vars_dict = template_vars or {}

            # Add default variables
            vars_dict.setdefault('content', content)
            vars_dict.setdefault('date', self._get_current_date())
            vars_dict.setdefault('timestamp', self._get_current_timestamp())

            rendered = template.render(**vars_dict)
            self.logger.info(f"Rendered content through HTML template: {template_name}")
            return rendered

        except Exception as e:
            self.logger.error(f"HTML template rendering failed: {e}")
            return content

    def _preprocess_markdown(self, text: str) -> str:
        """Preprocess markdown to handle common issues."""
        # Replace placeholder variables with empty strings
        text = re.sub(r'\{[a-zA-Z0-9_]+\}', '', text)

        # Handle f-strings that weren't evaluated
        text = re.sub(r'f"""(.*?)"""', r'\1', text, flags=re.DOTALL)
        text = re.sub(r"f'''(.*?)'''", r'\1', text, flags=re.DOTALL)

        # Remove triple backticks and language indicators
        text = re.sub(r'```[a-zA-Z]*\n', '', text)
        text = re.sub(r'```', '', text)

        # Fix heading issues (ensure space after #) - this should work correctly
        text = re.sub(r'(#+)([^ \n])', r'\1 \2', text)

        # Fix escaped newlines if any
        text = text.replace('\\n', '\n')

        # Clean up extra whitespace but preserve line structure
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            # Strip leading/trailing whitespace but preserve the line
            cleaned_line = line.strip()
            cleaned_lines.append(cleaned_line)

        return '\n'.join(cleaned_lines)

    def _markdown_to_html(self, markdown_text: str) -> str:
        """Convert markdown to HTML."""
        try:
            self.logger.debug(f"Converting markdown to HTML. Input preview: {markdown_text[:100]}...")

            html = markdown.markdown(
                markdown_text,
                extensions=['extra', 'codehilite', 'tables']  # Removed 'toc' to avoid issues
            )

            self.logger.debug(f"HTML conversion result preview: {html[:200]}...")
            return html

        except Exception as e:
            self.logger.error(f"Markdown conversion failed: {e}")
            # Fallback: wrap in paragraphs
            paragraphs = markdown_text.split('\n\n')
            html_paragraphs = [f'<p>{p.replace(chr(10), "<br>")}</p>' for p in paragraphs if p.strip()]
            fallback_html = '\n'.join(html_paragraphs)
            self.logger.debug(f"Using fallback HTML: {fallback_html[:200]}...")
            return fallback_html

    def _extract_slides_from_html(self, html_content: str, max_slides: int) -> List[Dict[str, Any]]:
        """Extract slides from HTML content based on headings."""
        soup = BeautifulSoup(html_content, 'html.parser')
        slides = []

        # Find all heading elements
        headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])

        if not headings:
            # If no headings, create a single slide with all content
            slides.append({
                'title': 'Presentation',
                'content': self._extract_content_elements(soup),
                'level': 1,
                'layout': 0  # Title slide for single content
            })
            return slides

        # Get all elements in the document to process sequentially
        all_elements = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'table', 'blockquote', 'div'])

        current_slide = None

        for element in all_elements:
            if len(slides) >= max_slides:
                self.logger.warning(f"Reached maximum slides limit ({max_slides}), stopping slide creation")
                break

            # If this is a heading, start a new slide
            if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                # Save the previous slide if it exists
                if current_slide is not None:
                    slides.append(current_slide)

                # Start new slide
                heading_level = int(element.name[1])
                heading_text = self._get_text_content(element)

                current_slide = {
                    'title': heading_text,
                    'level': heading_level,
                    'content': [],
                    'layout': 0 if (len(slides) == 0 and heading_level == 1) else 1
                }

                self.logger.debug(f"Starting new slide: {heading_text}")

            # If this is content and we have a current slide, add it
            elif element.name in ['p', 'ul', 'ol', 'table', 'blockquote', 'div'] and current_slide is not None:
                # Skip empty paragraphs
                content_text = self._get_text_content(element).strip()
                if content_text:
                    current_slide['content'].append(element)
                    self.logger.debug(f"Added content to slide '{current_slide['title']}': {content_text[:50]}...")

        # Don't forget the last slide
        if current_slide is not None:
            slides.append(current_slide)

        self.logger.info(f"Extracted {len(slides)} slides from HTML content")

        # Debug: Log slide information
        for i, slide in enumerate(slides):
            self.logger.debug(f"Slide {i+1}: '{slide['title']}' with {len(slide['content'])} content elements")

        return slides

    def _extract_content_elements(self, soup) -> List:
        """Extract content elements from soup."""
        content_elements = []

        # Get all content elements, but exclude headings since they become slide titles
        for element in soup.find_all(['p', 'ul', 'ol', 'table', 'blockquote', 'div']):
            # Skip if this div only contains headings
            if element.name == 'div':
                # Check if div contains only headings or is empty
                child_tags = [child.name for child in element.find_all() if hasattr(child, 'name')]
                if all(tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6'] for tag in child_tags):
                    continue

            content_elements.append(element)

        return content_elements

    def _create_presentation(self, template_path: Optional[str] = None) -> Presentation:
        """Create or load PowerPoint presentation."""
        if template_path:
            pptx_template = self._get_template_path(template_path)
            if pptx_template and pptx_template.exists():
                self.logger.info(f"Loading PowerPoint template: {pptx_template}")
                return Presentation(str(pptx_template))

        # Create new presentation
        return Presentation()

    def _create_slides(
        self,
        prs: Presentation,
        slides_data: List[Dict[str, Any]],
        slide_layout: int,
        title_styles: Optional[Dict[str, Any]],
        content_styles: Optional[Dict[str, Any]]
    ) -> None:
        """Create slides from extracted data."""
        for slide_data in slides_data:
            try:
                # Determine layout: use slide_data layout or default
                layout_idx = slide_data.get('layout', slide_layout)

                # Ensure layout index is valid
                if layout_idx >= len(prs.slide_layouts):
                    layout_idx = min(slide_layout, len(prs.slide_layouts) - 1)

                # Add slide
                slide_layout_obj = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(slide_layout_obj)

                # Add title
                if slide.shapes.title and slide_data['title']:
                    # first, remove any "#" at beginning of the title:
                    _title = slide_data['title'].lstrip('#').strip()
                    slide.shapes.title.text = _title
                    if title_styles:
                        self._apply_text_styles(slide.shapes.title, title_styles)

                # Add content (if there's a content placeholder and content exists)
                if slide_data['content'] and len(slide.shapes.placeholders) > 1:
                    content_placeholder = slide.shapes.placeholders[1]
                    self._add_slide_content(content_placeholder, slide_data['content'], content_styles)

                self.logger.debug(
                    f"Created slide: {slide_data['title']}"
                )

            except Exception as e:
                self.logger.error(f"Error creating slide '{slide_data.get('title', 'Unknown')}': {e}")
                continue

    def _add_slide_content(self, placeholder, content_elements: List, content_styles: Optional[Dict[str, Any]]) -> None:
        """Add content to a slide placeholder."""
        try:
            text_frame = placeholder.text_frame
            text_frame.clear()

            for element in content_elements:
                if element.name == 'p':
                    # Add paragraph
                    if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = self._get_text_content(element)
                    if content_styles:
                        self._apply_paragraph_styles(p, content_styles)

                elif element.name in ['ul', 'ol']:
                    # Add list items
                    for li in element.find_all('li', recursive=False):
                        p = text_frame.add_paragraph()
                        p.text = self._get_text_content(li)
                        p.level = 1  # Bullet point level
                        if content_styles:
                            self._apply_paragraph_styles(p, content_styles)

                elif element.name == 'table':
                    # Add table as formatted text
                    table_text = self._extract_table_text(element)
                    p = text_frame.add_paragraph()
                    p.text = table_text
                    if content_styles:
                        self._apply_paragraph_styles(p, content_styles)

                elif element.name == 'blockquote':
                    # Add blockquote
                    p = text_frame.add_paragraph()
                    p.text = f'"{self._get_text_content(element)}"'
                    if content_styles:
                        self._apply_paragraph_styles(p, content_styles)

        except Exception as e:
            self.logger.error(f"Error adding slide content: {e}")

    def _extract_table_text(self, table_element) -> str:
        """Extract text from table element."""
        rows = table_element.find_all('tr')
        table_lines = []

        for row in rows:
            cells = row.find_all(['td', 'th'])
            row_text = ' | '.join([self._get_text_content(cell) for cell in cells])
            table_lines.append(row_text)

        return '\n'.join(table_lines)

    def _get_text_content(self, element) -> str:
        """Extract clean text content from HTML element."""
        if isinstance(element, NavigableString):
            return str(element).strip()

        # For HTML elements, get the text content
        if hasattr(element, 'get_text'):
            # BeautifulSoup's get_text() method extracts clean text without HTML tags
            text = element.get_text(strip=True)
            return text

        # Fallback method for manual text extraction
        text_parts = []
        for content in element.contents:
            if isinstance(content, NavigableString):
                text_parts.append(str(content).strip())
            else:
                text_parts.append(self._get_text_content(content))

        result = ''.join(text_parts).strip()

        # Additional cleanup: remove any remaining markdown symbols
        # This shouldn't be necessary if markdown conversion worked correctly,
        # but it's a safety net
        result = re.sub(r'^#+\s*', '', result)  # Remove leading hashtags
        result = re.sub(r'\*\*([^*]+)\*\*', r'\1', result)  # Remove bold markers
        result = re.sub(r'\*([^*]+)\*', r'\1', result)  # Remove italic markers

        return result

    def _apply_text_styles(self, shape, styles: Dict[str, Any]) -> None:
        """Apply styles to a text shape."""
        if not shape.has_text_frame:
            return

        try:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                self._apply_paragraph_styles(paragraph, styles)
        except Exception as e:
            self.logger.error(f"Error applying text styles: {e}")

    def _apply_paragraph_styles(self, paragraph, styles: Dict[str, Any]) -> None:
        """Apply styles to a paragraph."""
        try:
            # Font styling
            if 'font_name' in styles:
                paragraph.font.name = styles['font_name']
            if 'font_size' in styles:
                paragraph.font.size = Pt(styles['font_size'])
            if 'bold' in styles:
                paragraph.font.bold = styles['bold']
            if 'italic' in styles:
                paragraph.font.italic = styles['italic']
            if 'font_color' in styles:
                # Convert hex color to RGB
                color_hex = styles['font_color'].lstrip('#')
                r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                paragraph.font.color.rgb = RGBColor(r, g, b)

            # Alignment
            if 'alignment' in styles:
                alignment_map = {
                    'left': PP_ALIGN.LEFT,
                    'center': PP_ALIGN.CENTER,
                    'right': PP_ALIGN.RIGHT,
                    'justify': PP_ALIGN.JUSTIFY
                }
                paragraph.alignment = alignment_map.get(styles['alignment'], PP_ALIGN.LEFT)

        except Exception as e:
            self.logger.error(f"Error applying paragraph styles: {e}")

    def debug_content_parsing(self, content: str) -> Dict[str, Any]:
        """
        Debug method to see how content is being parsed.

        Args:
            content: Input content to debug

        Returns:
            Dictionary with debug information
        """
        try:
            # Process the content the same way as in generation
            processed_content = self._preprocess_markdown(content)
            html_content = self._markdown_to_html(processed_content)

            # Parse with BeautifulSoup
            soup = BeautifulSoup(html_content, 'html.parser')
            headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])

            # Extract slide information
            slides_data = self._extract_slides_from_html(html_content, 50)

            debug_info = {
                "original_content_length": len(content),
                "original_content_preview": content[:300] + "..." if len(content) > 300 else content,
                "processed_content_preview": processed_content[:300] + "..." if len(processed_content) > 300 else processed_content,
                "html_content": html_content,  # Show full HTML to debug
                "headings_found": [
                    {
                        "tag": h.name,
                        "level": int(h.name[1]),
                        "raw_html": str(h),
                        "extracted_text": self._get_text_content(h),
                        "inner_text": h.get_text() if hasattr(h, 'get_text') else "N/A"
                    } for h in headings
                ],
                "slides_extracted": [
                    {
                        "title": slide['title'],
                        "level": slide['level'],
                        "layout": slide['layout'],
                        "content_count": len(slide['content']),
                        "content_preview": [
                            {
                                "tag": elem.name,
                                "text": self._get_text_content(elem)[:100]
                            } for elem in slide['content'][:3]
                        ]
                    } for slide in slides_data
                ],
                "total_slides": len(slides_data)
            }

            return debug_info

        except Exception as e:
            return {
                "error": str(e),
                "traceback": traceback.format_exc(),
                "message": "Debug parsing failed"
            }

    async def _generate_document_content(self, content: str, **kwargs) -> bytes:
        """
        Generate PowerPoint presentation content from input.

        Args:
            content: Input content (text, markdown, or HTML)
            **kwargs: Additional arguments from PowerPointArgs

        Returns:
            PowerPoint presentation as bytes
        """
        try:
            # Extract arguments
            template_name = kwargs.get('template_name')
            template_vars = kwargs.get('template_vars')
            pptx_template = kwargs.get('pptx_template')
            slide_layout = kwargs.get('slide_layout', 1)
            title_styles = kwargs.get('title_styles')
            content_styles = kwargs.get('content_styles')
            max_slides = kwargs.get('max_slides', 50)
            split_by_headings = kwargs.get('split_by_headings', True)

            # Process content through HTML template if provided
            processed_content = self._render_html_template(content, template_name, template_vars)

            # Preprocess markdown
            cleaned_content = self._preprocess_markdown(processed_content)

            # Convert to HTML
            html_content = self._markdown_to_html(cleaned_content)

            # Extract slides from HTML
            if split_by_headings:
                slides_data = self._extract_slides_from_html(html_content, max_slides)
            else:
                # Create single slide with all content
                soup = BeautifulSoup(html_content, 'html.parser')
                slides_data = [{
                    'title': 'Presentation',
                    'content': self._extract_content_elements(soup),
                    'level': 1,
                    'layout': 0
                }]

            self.logger.info(f"Generated {len(slides_data)} slides from content")

            # Create PowerPoint presentation
            prs = self._create_presentation(pptx_template)

            # Create slides
            self._create_slides(prs, slides_data, slide_layout, title_styles, content_styles)

            # Save presentation to bytes
            ppt_bytes = io.BytesIO()
            prs.save(ppt_bytes)
            ppt_bytes.seek(0)

            return ppt_bytes.getvalue()

        except Exception as e:
            self.logger.error(f"Error generating PowerPoint presentation: {e}")
            raise

    async def _execute(
        self,
        content: str,
        output_filename: Optional[str] = None,
        file_prefix: str = "presentation",
        output_dir: Optional[str] = None,
        overwrite_existing: bool = False,
        template_name: Optional[str] = None,
        template_vars: Optional[Dict[str, Any]] = None,
        pptx_template: Optional[str] = None,
        slide_layout: int = 1,
        title_styles: Optional[Dict[str, Any]] = None,
        content_styles: Optional[Dict[str, Any]] = None,
        max_slides: int = 50,
        split_by_headings: bool = True,
        **kwargs
    ) -> Dict[str, Any]:
        """
        Execute PowerPoint presentation generation (AbstractTool interface).

        Args:
            content: Content to convert to PowerPoint presentation
            output_filename: Custom filename (without extension)
            file_prefix: Prefix for auto-generated filenames
            output_dir: Custom output directory
            overwrite_existing: Whether to overwrite existing files
            template_name: HTML template name for content processing
            template_vars: Variables for HTML template
            pptx_template: PowerPoint template file path
            slide_layout: Default slide layout index
            title_styles: Styles for slide titles
            content_styles: Styles for slide content
            max_slides: Maximum number of slides to generate
            split_by_headings: Whether to split by headings
            **kwargs: Additional arguments

        Returns:
            Dictionary with presentation generation results
        """
        try:
            self.logger.info(f"Starting PowerPoint generation with {len(content)} characters of content")

            # Use the safe document creation workflow
            result = await self._create_document_safely(
                content=content,
                output_filename=output_filename,
                file_prefix=file_prefix,
                output_dir=output_dir,
                overwrite_existing=overwrite_existing,
                extension="pptx",
                template_name=template_name,
                template_vars=template_vars,
                pptx_template=pptx_template,
                slide_layout=slide_layout,
                title_styles=title_styles,
                content_styles=content_styles,
                max_slides=max_slides,
                split_by_headings=split_by_headings
            )

            if result['status'] == 'success':
                # Add presentation-specific metadata
                result['presentation_info'] = {
                    'max_slides_limit': max_slides,
                    'split_by_headings': split_by_headings,
                    'slide_layout_used': slide_layout
                }

                self.logger.debug(
                    f"PowerPoint presentation created successfully: {result['metadata']['filename']}"
                )

            return result

        except Exception as e:
            self.logger.error(f"Error in PowerPoint generation: {e}")
            raise
